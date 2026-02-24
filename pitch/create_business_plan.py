from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Create presentation (same size as template: 10x5.62)
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
BLUE = RGBColor(0, 102, 255)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_GRAY = RGBColor(50, 50, 50)
BLACK = RGBColor(0, 0, 0)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Blue header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_items, slide_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Blue header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Slide number
    if slide_num:
        num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(8.5), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

def add_two_column_slide(title, left_title, left_items, right_title, right_items, slide_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Blue header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Slide number
    if slide_num:
        num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(0.5), Inches(0.4))
        p = num_box.text_frame.paragraphs[0]
        p.text = slide_num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(8.5), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.2), Inches(0.4))
    p = left_header.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.2), Inches(3.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Right column
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.2), Inches(0.4))
    p = right_header.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.2), Inches(3.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    return slide

# ============ BUILD THE SLIDES ============

# SLIDE 1: Title
add_title_slide("MOTRA", "Autonomy, Maintained\n\nBusiness Plan")

# SLIDE 2: Team
add_content_slide("Team", [
    "ADONIS WILLIAMS — Founder & CEO",
    "",
    "Current Role: Deputy Functional Chief Engineer at Boeing",
    "Experience: Engineering Manager, Boeing Research & Technology",
    "Education: Missouri University of Science and Technology",
    "Location: Seattle, Washington — epicenter of tech and mobility innovation",
    "",
    "Why Adonis:",
    "• 10+ years leading complex engineering systems at scale",
    "• Deep expertise in fleet operations and infrastructure",
    "• Track record of building and shipping products used by millions",
    "• Connected to Seattle's tech and mobility ecosystem"
], "1")

# SLIDE 3: Team Values
add_two_column_slide("Team Values & Norms", 
    "Values",
    ["Engineering Excellence", "Scalability First", "Operator Mindset", "Data-Driven Decisions", "Customer Obsession"],
    "Norms",
    ["Own outcomes end-to-end", "Test assumptions rapidly", "Measure everything", "Move fast once decided", "Weekly iteration cycles"],
    "2"
)

# SLIDE 4: Define the Need
add_two_column_slide("Define the Need: Overarching Hypothesis",
    "USER",
    ["Autonomous Vehicle Fleet Operators", "Companies like Waymo, Cruise, Zoox, Tesla", "Fleet Operations Managers", "VP of Operations / Head of Fleet"],
    "PAIN POINT",
    ["Who cleans robotaxis between rides?", "No driver to notice spills, trash, wear", "Depot-based cleaning = downtime = lost revenue", "Sensors must stay clean (safety-critical)", "Current solutions don't scale to millions of rides"],
    "3"
)

# SLIDE 5: Hypothesis
add_content_slide("Our Hypothesis", [
    "AV fleet operators will pay for on-demand, mobile cleaning and maintenance services",
    "that come to their vehicles — eliminating depot trips and maximizing uptime.",
    "",
    "SOLUTION:",
    "MOTRA deploys a gig-powered network of mobile technicians to clean and service",
    "autonomous vehicles on-location, on-demand.",
    "",
    "KEY INSIGHT:",
    "Waymo currently hand-washes every vehicle at centralized depots.",
    "This approach fundamentally doesn't scale to millions of rides.",
    "",
    "MOTRA is AWS for autonomous fleet maintenance —",
    "invisible, essential, everywhere."
], "4")

# SLIDE 6: User Journey
add_content_slide("User Journey", [
    "1. INSPIRATION",
    "   Fleet operator realizes cleaning ops are becoming a bottleneck",
    "   Vehicles spending too much time in depot, not earning",
    "",
    "2. DISCOVERY",
    "   Searches for fleet cleaning solutions, finds MOTRA",
    "   Sees: mobile, on-demand, gig-powered, AV-specialized",
    "",
    "3. DECISION",
    "   Runs pilot with 50 vehicles in one city",
    "   Compares: cost per clean, vehicle uptime, quality scores",
    "",
    "4. EXECUTION",
    "   Integrates MOTRA API with fleet management system",
    "   Automated dispatch after X rides or on-demand triggers",
    "",
    "5. EXPANSION",
    "   Rolls out to all markets, becomes preferred vendor"
], "5")

# SLIDE 7: Facts & Assumptions
add_two_column_slide("Facts vs. Hypotheses",
    "VALIDATED FACTS",
    ["Waymo: 400K rides/week, targeting 1M/week by 2026", 
     "Fleet: 2,500 vehicles → 'tens of thousands' planned",
     "Waymo hand-washes at depots (can't use car washes)",
     "15 million rides delivered in 2025 alone",
     "Tesla Robotaxi, Zoox, Cruise all scaling 2026-2027"],
    "HYPOTHESES TO TEST",
    ["Fleet operators want to outsource cleaning",
     "Mobile service is faster than depot round-trips",
     "$15/service price point is acceptable",
     "Gig workers can be trained to AV standards",
     "Quality can be maintained at scale"],
    "6"
)

# SLIDE 8: Market Size
add_content_slide("Market Size", [
    "TAM (Total Addressable Market) — 2032",
    "500,000 AVs × 2 services/day × $15/service × 365 days = $5.5 BILLION",
    "",
    "SAM (Serviceable Addressable Market) — 2028",
    "50,000 AVs × 2 services/day × $15/service × 365 days = $547 MILLION",
    "",
    "SOM (Serviceable Obtainable Market) — Year 1",
    "10% of one operator (250 vehicles) in one city = $2 MILLION",
    "",
    "GROWTH DRIVERS:",
    "• Waymo targeting 1M rides/week by end of 2026",
    "• Tesla Robotaxi launching 2026",
    "• Zoox (Amazon) expanding",
    "• EV fleet services (Amazon vans, FedEx, UPS) — adjacent market"
], "7")

# SLIDE 9: Competition
add_two_column_slide("Competitive Landscape",
    "CURRENT ALTERNATIVES",
    ["AV In-House Ops: High fixed cost, doesn't scale",
     "Traditional Fleet Services: Not mobile, not AV-specialized",
     "Car Washes: Damages sensors, no interior service",
     "Mobile Detailing: Not scaled, not fleet-focused",
     "Gig Platforms: No AV specialization or training"],
    "MOTRA'S ADVANTAGES",
    ["AV-Specialized: Sensors, EVs, safety protocols",
     "Mobile-First: Go to the vehicle, not vice versa",
     "Gig-Powered: Variable cost, instant scaling",
     "Tech-Enabled: API integration with fleet systems",
     "First Mover: No scaled competitor exists"],
    "8"
)

# SLIDE 10: Traction
add_content_slide("Traction To Date", [
    "COMPLETED:",
    "• Market research: Validated Waymo's depot-based cleaning bottleneck",
    "• TAM/SAM/SOM analysis complete",
    "• Competitive landscape mapped — blue ocean confirmed",
    "• Landing page live: tidytails.github.io/motra-landing",
    "• Pitch materials ready",
    "",
    "IN PROGRESS:",
    "• Outreach to Waymo Fleet Operations contacts",
    "• Identifying pilot market (Phoenix or SF)",
    "• Defining MVP service scope",
    "",
    "NEXT MILESTONES:",
    "• First customer discovery call with AV operator",
    "• Letter of intent from pilot partner"
], "9")

# SLIDE 11: North Star
add_content_slide("North Star Vision", [
    "",
    "\"Every autonomous vehicle in the world is serviced by",
    "MOTRA's invisible infrastructure layer —",
    "clean, maintained, and always ready.\"",
    "",
    "",
    "We become the AWS of autonomous fleet maintenance:",
    "",
    "• Present in every market where AVs operate",
    "• Integrated into every fleet management system",
    "• The default answer when operators think 'fleet care'",
    "• Enabling the autonomous future by handling what humans used to do"
], "10")

# SLIDE 12: Business Model Canvas Overview
add_content_slide("Business Model Canvas", [
    "KEY PARTNERS                          KEY ACTIVITIES                    VALUE PROPOSITION",
    "• AV manufacturers                    • Platform development            • Reduce fleet downtime",
    "• Fleet management systems            • Tech recruitment/training       • Variable cost model",
    "• Gig worker platforms                • Quality assurance               • 24/7 availability",
    "• Insurance providers                 • Enterprise sales                • AV-specialized service",
    "",
    "KEY RESOURCES                         CUSTOMER RELATIONSHIPS            CHANNELS",
    "• Tech platform                       • Dedicated account managers      • Direct enterprise sales",
    "• Trained technician network          • API integration support         • Industry conferences",
    "• AV-specific training program        • Performance dashboards          • Partnerships",
    "",
    "COST STRUCTURE                                    REVENUE STREAMS",
    "• Platform development & maintenance              • Per-service fees ($12-75)",
    "• Technician payouts (65-70% of service fee)      • Enterprise contracts",
    "• Marketing & sales                               • Premium/emergency services",
    "• Insurance & compliance                          • Data analytics (future)"
], "11")

# SLIDE 13: Value Proposition
add_content_slide("Value Proposition", [
    "FOR AV FLEET OPERATORS:",
    "",
    "• REDUCE DOWNTIME: Vehicles stay in service zones, no depot trips",
    "• VARIABLE COSTS: Pay per service, not per facility",
    "• SCALE INSTANTLY: Expand to new markets without building infrastructure",
    "• QUALITY ASSURED: Certified techs, standardized processes, tracked metrics",
    "• API INTEGRATION: Automated dispatch based on ride data",
    "",
    "VS. CURRENT ALTERNATIVES:",
    "",
    "• 50% cost reduction vs. in-house depot operations",
    "• 30% more vehicle uptime",
    "• Zero facility buildout in new markets",
    "• Real-time quality tracking (in-house is often unmeasured)"
], "12")

# SLIDE 14: Customer Segments
add_content_slide("Customer Segments", [
    "PRIMARY: Autonomous Vehicle Fleet Operators",
    "• Waymo (Alphabet) — 2,500+ vehicles, 400K rides/week",
    "• Cruise (GM) — Rebuilding after 2023 pause",
    "• Zoox (Amazon) — Testing in SF, Las Vegas",
    "• Tesla Robotaxi — Launching 2026",
    "",
    "SECONDARY: Electric Fleet Operators",
    "• Amazon Delivery Vans (100,000+ vehicles)",
    "• FedEx/UPS Electric Fleets",
    "• Corporate EV Fleets",
    "",
    "BUYER PERSONAS:",
    "• VP of Operations / Head of Fleet",
    "• Fleet Operations Manager",
    "• Facilities/Maintenance Director",
    "",
    "DECISION CRITERIA: Cost per clean, uptime improvement, quality scores, scalability"
], "13")

# SLIDE 15: Revenue Model
add_content_slide("Revenue Streams", [
    "PER-SERVICE PRICING:",
    "• Quick Clean (5-10 min): $12-18",
    "• Deep Clean (30-60 min): $45-75",
    "• Light Maintenance: $25-50",
    "• Emergency Response: $75-150 (premium)",
    "",
    "UNIT ECONOMICS:",
    "• Average service price: $15",
    "• Technician payout: $10-11 (65-70%)",
    "• Platform margin: $4-5 (27-33%)",
    "",
    "AT SCALE:",
    "• 1,000 services/day = $15K revenue, $4-5K margin",
    "• 10,000 services/day = $150K revenue, $40-50K margin",
    "",
    "ENTERPRISE CONTRACTS:",
    "• Volume discounts for committed minimums",
    "• SLA guarantees (response time, quality scores)"
], "14")

# SLIDE 16: Key Activities
add_content_slide("Key Activities", [
    "PLATFORM DEVELOPMENT:",
    "• Fleet operator dashboard (scheduling, tracking, analytics)",
    "• Technician mobile app (dispatch, checklists, earnings)",
    "• API for fleet management system integration",
    "",
    "NETWORK OPERATIONS:",
    "• Technician recruitment and onboarding",
    "• AV-specific training and certification",
    "• Quality assurance and performance management",
    "• Equipment kit distribution",
    "",
    "ENTERPRISE SALES:",
    "• Outreach to AV fleet operators",
    "• Pilot program management",
    "• Account expansion and retention"
], "15")

# SLIDE 17: Key Resources
add_content_slide("Key Resources", [
    "TECHNOLOGY:",
    "• Dispatch and routing platform",
    "• Fleet operator dashboard",
    "• Technician mobile application",
    "• API integration layer",
    "",
    "HUMAN:",
    "• Engineering team (platform development)",
    "• Operations team (network management)",
    "• Sales team (enterprise relationships)",
    "• Trained technician network",
    "",
    "INTELLECTUAL:",
    "• AV-specific training curriculum",
    "• Quality standards and checklists",
    "• Operational playbooks for each market",
    "",
    "FINANCIAL:",
    "• Seed funding for 18-month runway"
], "16")

# SLIDE 18: Key Partners
add_content_slide("Key Partners", [
    "AV ECOSYSTEM:",
    "• Waymo, Cruise, Zoox, Tesla (customers)",
    "• AV industry associations",
    "• Fleet management software providers",
    "",
    "GIG ECONOMY:",
    "• Existing gig platforms (recruitment channel)",
    "• Detailing/cleaning training providers",
    "• Equipment suppliers",
    "",
    "SUPPORTING:",
    "• Insurance providers (commercial auto, liability)",
    "• Background check services",
    "• Payment processing (Stripe, etc.)",
    "",
    "STRATEGIC:",
    "• EV charging networks (co-location opportunities)",
    "• Parking/logistics providers"
], "17")

# SLIDE 19: Cost Structure
add_content_slide("Cost Structure", [
    "VARIABLE COSTS (scale with volume):",
    "• Technician payouts: 65-70% of service revenue",
    "• Payment processing: 2-3%",
    "• Insurance per service",
    "",
    "FIXED COSTS:",
    "• Platform development & maintenance: $15-20K/month",
    "• Core team salaries: $40-60K/month",
    "• Marketing & sales: $10-15K/month",
    "• Insurance (business): $2-5K/month",
    "• Legal & compliance: $3-5K/month",
    "",
    "STARTUP COSTS:",
    "• MVP development: $100-150K",
    "• Initial market launch: $50K",
    "• Working capital: $100K",
    "",
    "TOTAL SEED NEED: $1.5M for 18-month runway"
], "18")

# SLIDE 20: Financial Projections
add_content_slide("Financial Projections", [
    "YEAR 1:",
    "• Markets: 1 (Phoenix or SF)",
    "• Vehicles served: 500",
    "• Revenue: $1.2M",
    "• Net: -$300K (investment phase)",
    "",
    "YEAR 2:",
    "• Markets: 3",
    "• Vehicles served: 3,000",
    "• Revenue: $5.5M",
    "• Net: +$800K",
    "",
    "YEAR 3:",
    "• Markets: 7",
    "• Vehicles served: 10,000",
    "• Revenue: $15M",
    "• Net: +$3.5M",
    "",
    "PATH TO $10M ARR: 36 months"
], "19")

# SLIDE 21: The Ask
add_content_slide("The Ask", [
    "RAISING: $1.5M Seed Round",
    "",
    "USE OF FUNDS:",
    "• Product Development (40%): Platform, apps, API integrations",
    "• Operations (30%): Tech recruitment, training, equipment",
    "• Sales (20%): Enterprise BD, pilot programs",
    "• G&A (10%): Legal, insurance, admin",
    "",
    "MILESTONES TO SERIES A:",
    "• 1 enterprise contract with major AV operator",
    "• 1,000+ services/month",
    "• Positive unit economics proven",
    "• Expansion to 2nd market",
    "",
    "RUNWAY: 18 months to Series A metrics"
], "20")

# SLIDE 22: Closing
add_title_slide("MOTRA", "Autonomy, Maintained\n\nAdonis Williams | Seattle, WA\ntidytails.github.io/motra-landing")

# Save
output_path = "/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx"
prs.save(output_path)
print(f"Saved comprehensive business plan to {output_path}")
print(f"Total slides: {len(prs.slides)}")
