#!/usr/bin/env python3
"""
MOTRA Pitch Deck v3 - PowerPoint Generator (Fixed Alignment)
Updated March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RgbColor(0, 102, 255)
DARK = RgbColor(26, 26, 26)
GRAY = RgbColor(102, 102, 102)
WHITE = RgbColor(255, 255, 255)
RED = RgbColor(220, 38, 38)
GREEN = RgbColor(22, 163, 74)

def add_text_box(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box

def add_stat_box(slide, left, top, width, value, label, value_size=32):
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(0, 40, 80)
    box.line.color.rgb = BLUE
    
    # Value - centered in box
    val_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.15), Inches(width), Inches(0.6))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(value_size)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    lbl_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.8), Inches(width), Inches(0.4))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_card(slide, left, top, width, height, title, content_lines, title_color=WHITE):
    # Card background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(30, 30, 30)
    box.line.color.rgb = RgbColor(60, 60, 60)
    
    # Title
    add_text_box(slide, left + 0.15, top + 0.15, width - 0.3, 0.4, title, size=14, color=title_color, bold=True)
    
    # Content lines
    y = top + 0.55
    for line in content_lines:
        add_text_box(slide, left + 0.15, y, width - 0.3, 0.35, "• " + line, size=11, color=GRAY)
        y += 0.35

# === SLIDE 0: Title ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 2.8, 12.333, 1, "MOTRA", size=72, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4.2, 12.333, 0.5, "AUTONOMY, MAINTAINED", size=24, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 6.5, 12.333, 0.4, "Seed Round — March 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# === SLIDE 1: Company Purpose ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "01 / COMPANY PURPOSE", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "The infrastructure layer for AV fleet care.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.9, 12, 0.8, "MOTRA is the infrastructure layer for autonomous vehicle fleet care.", size=24, color=WHITE)
add_text_box(slide, 0.5, 3.0, 12, 1, "We deploy a gig-powered network of mobile technicians to clean and service robotaxis on-location, on-demand — keeping fleets running 24/7.", size=16, color=GRAY)

# Highlight box
hbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12), Inches(1))
hbox.fill.solid()
hbox.fill.fore_color.rgb = RgbColor(0, 40, 80)
hbox.line.color.rgb = BLUE
add_text_box(slide, 0.7, 4.75, 11.5, 0.5, "Think DoorDash for fleet maintenance — invisible, essential, everywhere.", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# === SLIDE 2: Problem ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "02 / THE PROBLEM", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "The maintenance bottleneck.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.8, 12, 0.6, "Robotaxis run 15+ hours a day. Their maintenance infrastructure doesn't.", size=20, color=WHITE)

# Pain points
points = [
    ("1. No driver = no eyes", "Nobody notices trash, spills, or wear"),
    ("2. Depot-dependent", "Vehicles must return for basic cleaning"),
    ("3. Downtime = lost revenue", "Every minute in depot is a ride not taken"),
    ("4. Sensors are critical", "Dirty LiDAR/cameras = degraded driving")
]
y = 2.6
for title, desc in points:
    add_text_box(slide, 0.5, y, 5, 0.35, title, size=14, color=WHITE, bold=True)
    add_text_box(slide, 0.5, y + 0.35, 5, 0.3, desc, size=11, color=GRAY)
    y += 0.75

# Current reality box
rbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.6), Inches(6), Inches(2.8))
rbox.fill.solid()
rbox.fill.fore_color.rgb = RgbColor(50, 30, 30)
rbox.line.color.rgb = RED
add_text_box(slide, 6.7, 2.8, 5.5, 0.4, "Current Reality", size=16, color=RED, bold=True)
add_text_box(slide, 6.7, 3.3, 5.5, 0.5, "Waymo uses 1,000+ Transdev operators at centralized depots.", size=14, color=WHITE)
add_text_box(slide, 6.7, 4.0, 5.5, 0.8, "That doesn't scale to 1 million rides/week.", size=18, color=WHITE, bold=True)

# === SLIDE 3: Solution ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "03 / THE SOLUTION", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Mobile fleet care on demand.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.8, 12, 0.5, "MOTRA deploys certified technicians directly to vehicles — wherever they are.", size=16, color=GRAY)

# Service cards
services = [
    ("Quick Clean", "$12-18", "5-10 min", "Wipe-down, trash, odor"),
    ("Deep Clean", "$45-75", "30-60 min", "Full detail, sensors, exterior"),
    ("Emergency", "$75-150", "On-demand", "Spills, incidents, urgent")
]
x = 0.5
for title, price, time, desc in services:
    cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.6), Inches(3.9), Inches(2.2))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RgbColor(30, 30, 30)
    cbox.line.color.rgb = RgbColor(60, 60, 60)
    add_text_box(slide, x + 0.2, 2.8, 3.5, 0.4, title, size=18, color=WHITE, bold=True)
    add_text_box(slide, x + 0.2, 3.25, 3.5, 0.5, price, size=28, color=BLUE, bold=True)
    add_text_box(slide, x + 0.2, 3.8, 3.5, 0.3, time, size=12, color=GRAY)
    add_text_box(slide, x + 0.2, 4.15, 3.5, 0.4, desc, size=11, color=GRAY)
    x += 4.1

# Benefits bar
bbar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12), Inches(0.8))
bbar.fill.solid()
bbar.fill.fore_color.rgb = RgbColor(0, 40, 80)
bbar.line.color.rgb = BLUE
add_text_box(slide, 0.7, 5.4, 11.5, 0.4, "Vehicles stay in zone  |  Variable cost  |  Scales instantly  |  24/7 availability", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# === SLIDE 4: Why Now ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "04 / WHY NOW", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Escape velocity.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "The robotaxi industry just hit an inflection point.", size=18, color=GRAY)

# Stats
stats = [("400K", "Rides/Week"), ("10", "US Markets"), ("$16B", "Fresh Funding"), ("$126B", "Valuation")]
x = 0.5
for val, lbl in stats:
    add_stat_box(slide, x, 2.4, 2.9, val, lbl)
    x += 3.1

# Table header
add_text_box(slide, 0.5, 4.0, 2.5, 0.35, "Metric", size=11, color=BLUE, bold=True)
add_text_box(slide, 3.2, 4.0, 3.5, 0.35, "Now (Feb 2026)", size=11, color=BLUE, bold=True)
add_text_box(slide, 7, 4.0, 5, 0.35, "EOY 2026 Target", size=11, color=BLUE, bold=True)

# Table rows
rows = [
    ("Weekly Rides", "400,000", "1,000,000"),
    ("Markets", "10 cities", "20+ cities (London, Tokyo)"),
    ("Fleet Size", "2,500 vehicles", "4,500+ vehicles"),
    ("Price War", "Waymo $8.17 vs Uber $17.25", "Volume exploding")
]
y = 4.4
for metric, now, target in rows:
    add_text_box(slide, 0.5, y, 2.5, 0.35, metric, size=12, color=WHITE)
    add_text_box(slide, 3.2, y, 3.5, 0.35, now, size=12, color=WHITE, bold=True)
    add_text_box(slide, 7, y, 5, 0.35, target, size=12, color=WHITE)
    y += 0.4

add_text_box(slide, 0.5, 6.2, 12, 0.3, "Sources: TechXplore Feb 2026, Fox News Feb 2026, Bloomberg", size=9, color=GRAY)

# === SLIDE 5: Market Size ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "05 / MARKET SIZE", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "A $147B opportunity.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "Robotaxi market: 99% CAGR through 2033", size=18, color=GRAY)

# Big stats
add_stat_box(slide, 0.5, 2.4, 3.8, "$2.7B", "2025 TAM", value_size=36)
add_stat_box(slide, 4.6, 2.4, 3.8, "$147B", "2033 TAM", value_size=36)
add_stat_box(slide, 8.7, 2.4, 3.8, "99%", "CAGR", value_size=36)

# MOTRA serviceable market
add_text_box(slide, 0.5, 4.1, 12, 0.4, "MOTRA's Serviceable Market", size=16, color=BLUE, bold=True)
rows = [
    ("2028", "$547M", "100K robotaxis × ~$105/week"),
    ("2030", "$2.2B", "200K robotaxis"),
    ("2032", "$5.5B", "500K robotaxis + commercial EVs")
]
y = 4.6
for year, tam, note in rows:
    add_text_box(slide, 0.5, y, 1.2, 0.35, year, size=14, color=WHITE, bold=True)
    add_text_box(slide, 2, y, 1.8, 0.35, tam, size=14, color=BLUE, bold=True)
    add_text_box(slide, 4.2, y, 8, 0.35, note, size=12, color=GRAY)
    y += 0.4

add_text_box(slide, 0.5, 6.2, 12, 0.3, "Source: Grand View Research, MOTRA analysis", size=9, color=GRAY)

# === SLIDE 6: Competition ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "06 / COMPETITION", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Our competitive edge.", size=32, color=WHITE, bold=True)

# Table header
add_text_box(slide, 0.5, 1.9, 3, 0.35, "Player", size=11, color=BLUE, bold=True)
add_text_box(slide, 4, 1.9, 4, 0.35, "Approach", size=11, color=BLUE, bold=True)
add_text_box(slide, 8.5, 1.9, 4, 0.35, "Limitation", size=11, color=BLUE, bold=True)

rows = [
    ("Uber AV Services", "Depot-based cleaning + charging", "Competitor to Waymo"),
    ("Transdev", "1,000+ operators at depots", "Depot-bound, not mobile"),
    ("Tesla Robots", "Automated cleaning (patents)", "Years away, Tesla-only"),
    ("In-House Ops", "AV companies DIY", "High fixed cost")
]
y = 2.4
for player, approach, limitation in rows:
    add_text_box(slide, 0.5, y, 3.3, 0.4, player, size=13, color=WHITE, bold=True)
    add_text_box(slide, 4, y, 4, 0.4, approach, size=11, color=GRAY)
    add_text_box(slide, 8.5, y, 4, 0.4, limitation, size=11, color=RED)
    y += 0.5

# MOTRA advantage
abox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12), Inches(1))
abox.fill.solid()
abox.fill.fore_color.rgb = RgbColor(20, 50, 30)
abox.line.color.rgb = GREEN
add_text_box(slide, 0.7, 5.05, 11.5, 0.5, "MOTRA Advantage: Neutral (Switzerland) · Mobile-first · 100% focused on AV fleet care", size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# === SLIDE 7: Why Not In-House ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "07 / WHY NOT IN-HOUSE", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "The outsourcing thesis.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "Why won't Waymo just do this themselves?", size=18, color=GRAY)

# Three reasons
reasons = [
    ("1. Already Outsource", "Waymo uses Transdev, Moove, Avis, TechForce, Amerit, Terawatt today.", "Build vs. buy = they buy."),
    ("2. Not Core Competency", "Waymo's edge is autonomous driving AI.", "Every $ on cleaning = $ not on AI."),
    ("3. Math Doesn't Work", "20+ cities = thousands of employees, HR, facilities.", "Gig model = variable cost.")
]
x = 0.5
for title, line1, line2 in reasons:
    cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(3.9), Inches(2.4))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RgbColor(30, 30, 30)
    cbox.line.color.rgb = BLUE
    add_text_box(slide, x + 0.15, 2.55, 3.6, 0.4, title, size=14, color=BLUE, bold=True)
    add_text_box(slide, x + 0.15, 3.0, 3.6, 0.8, line1, size=11, color=WHITE)
    add_text_box(slide, x + 0.15, 3.8, 3.6, 0.5, line2, size=11, color=GREEN)
    x += 4.1

# Proof box
pbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12), Inches(0.9))
pbox.fill.solid()
pbox.fill.fore_color.rgb = RgbColor(0, 40, 80)
pbox.line.color.rgb = BLUE
add_text_box(slide, 0.7, 5.4, 11.5, 0.5, "Uber just validated this — if big players could do it in-house, they wouldn't build outsourced solutions.", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# === SLIDE 8: Moat ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "08 / DEFENSIBILITY", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Our moat.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "Three layers of defensibility.", size=18, color=GRAY)

moats = [
    ("Year 1-2: Network Effects", ["First AV gig network per city", "Hyperlocal network effects", "More techs → faster response → more volume"]),
    ("Year 2-3: Data Advantage", ["Every service = operational data", "Predictive maintenance insights", "DoorDash: 30% of eng in data roles"]),
    ("Year 3+: Integration Lock-In", ["API embedded in fleet systems", "Part of operator workflows", "5% retention = 25-95% profit boost"])
]
x = 0.5
for title, points in moats:
    cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(3.9), Inches(2.8))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RgbColor(30, 30, 30)
    cbox.line.color.rgb = BLUE
    add_text_box(slide, x + 0.15, 2.55, 3.6, 0.45, title, size=12, color=BLUE, bold=True)
    y = 3.1
    for pt in points:
        add_text_box(slide, x + 0.15, y, 3.6, 0.4, "• " + pt, size=10, color=WHITE)
        y += 0.4
    x += 4.1

add_text_box(slide, 0.5, 5.5, 12, 0.4, "The Uber analogy: Uber's moat isn't the app — it's the network of drivers in every city.", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# === SLIDE 9: Team ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "09 / TEAM", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Leadership.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "15 years at Boeing. 2 patents. Built to scale.", size=18, color=GRAY)

# Team card
tbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.4), Inches(6), Inches(3.2))
tbox.fill.solid()
tbox.fill.fore_color.rgb = RgbColor(30, 30, 30)
tbox.line.color.rgb = BLUE

add_text_box(slide, 0.7, 2.6, 5.5, 0.5, "Adonis Williams — Founder & CEO", size=16, color=WHITE, bold=True)

details = [
    ("Role:", "Deputy Functional Chief Engineer, Boeing"),
    ("Tenure:", "15 years at Boeing"),
    ("Patents:", "2 (cybersecurity, blockchain comms)"),
    ("Education:", "MS Systems Engineering, Missouri S&T")
]
y = 3.2
for label, value in details:
    add_text_box(slide, 0.7, y, 1.5, 0.35, label, size=10, color=GRAY)
    add_text_box(slide, 2.3, y, 4, 0.35, value, size=11, color=WHITE)
    y += 0.4

# Why Adonis box
wbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.4), Inches(5.5), Inches(3.2))
wbox.fill.solid()
wbox.fill.fore_color.rgb = RgbColor(20, 25, 35)
wbox.line.fill.background()

add_text_box(slide, 7.2, 2.6, 5, 0.4, "Why Adonis Wins This", size=14, color=BLUE, bold=True)
add_text_box(slide, 7.2, 3.1, 5, 1.2, "MOTRA is a systems problem: coordinating a distributed workforce to service safety-critical AVs.\n\nThat's exactly what he's done at Boeing for 15 years.", size=11, color=WHITE)
add_text_box(slide, 7.2, 4.4, 5, 1, "✓ Complex operations at scale\n✓ Distributed teams\n✓ Safety-critical systems\n✓ 2 patents", size=11, color=GREEN)

# === SLIDE 10: Risk Mitigation ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "10 / RISK MITIGATION", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "De-risked.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "What if the AV market slows down?", size=18, color=GRAY)

# Risk/Mitigation boxes
rbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.4), Inches(5.8), Inches(1.4))
rbox.fill.solid()
rbox.fill.fore_color.rgb = RgbColor(50, 30, 30)
rbox.line.color.rgb = RED
add_text_box(slide, 0.7, 2.55, 5.4, 0.35, "The Risk", size=14, color=RED, bold=True)
add_text_box(slide, 0.7, 2.95, 5.4, 0.7, "Robotaxi adoption could hit regulatory delays, technical hurdles, or economic headwinds.", size=11, color=WHITE)

mbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.4), Inches(5.8), Inches(1.4))
mbox.fill.solid()
mbox.fill.fore_color.rgb = RgbColor(20, 50, 30)
mbox.line.color.rgb = GREEN
add_text_box(slide, 7, 2.55, 5.4, 0.35, "Adjacent Markets", size=14, color=GREEN, bold=True)
add_text_box(slide, 7, 2.95, 5.4, 0.7, "Even if AV slows, 250K+ commercial EVs need service.", size=11, color=WHITE)

# Fleet table
add_text_box(slide, 0.5, 4.1, 3, 0.35, "Fleet", size=11, color=BLUE, bold=True)
add_text_box(slide, 4, 4.1, 3, 0.35, "Current", size=11, color=BLUE, bold=True)
add_text_box(slide, 8, 4.1, 4, 0.35, "Target", size=11, color=BLUE, bold=True)

fleets = [
    ("Amazon (Rivian)", "30,000 EVs", "100,000 by 2030"),
    ("FedEx", "200,000+ vehicles", "All-electric by 2040"),
    ("UPS", "18,300+ alt-fuel", "10,000 more ordered")
]
y = 4.5
for fleet, current, target in fleets:
    add_text_box(slide, 0.5, y, 3.3, 0.35, fleet, size=12, color=WHITE, bold=True)
    add_text_box(slide, 4, y, 3.5, 0.35, current, size=11, color=WHITE)
    add_text_box(slide, 8, y, 4, 0.35, target, size=11, color=WHITE)
    y += 0.4

add_text_box(slide, 0.5, 5.9, 12, 0.4, "MOTRA has multiple paths to scale. Robotaxis are the beachhead — commercial EVs are the expansion.", size=14, color=GREEN, align=PP_ALIGN.CENTER)

# === SLIDE 11: GTM ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "11 / GO-TO-MARKET", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "First customers.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "The DoorDash playbook: do things that don't scale.", size=18, color=GRAY)

phases = [
    ("Phase 1: Manual Pilot", "Month 1-3", ["ONE operator, ONE city", "Free 30-day, 50 vehicles", "Text dispatch, Google Sheets", "We ARE the operators"]),
    ("Phase 2: Paid Contract", "Month 4-6", ["Convert to paid", "Price 20% below Transdev", "6-month commitment", "Build reference customer"]),
    ("Phase 3: Scale", "Month 7-12", ["Add second city", "Launch MVP app", "First customer as reference", "Target $100K MRR"])
]
x = 0.5
for title, timeline, points in phases:
    cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(3.9), Inches(2.8))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RgbColor(30, 30, 30)
    cbox.line.color.rgb = BLUE
    add_text_box(slide, x + 0.15, 2.55, 3.6, 0.4, title, size=13, color=BLUE, bold=True)
    add_text_box(slide, x + 0.15, 2.95, 3.6, 0.3, timeline, size=10, color=GRAY)
    y = 3.3
    for pt in points:
        add_text_box(slide, x + 0.15, y, 3.6, 0.35, "• " + pt, size=10, color=WHITE)
        y += 0.35
    x += 4.1

add_text_box(slide, 0.5, 5.5, 12, 0.4, "Warm intros = 3-5x better close rate. We're asking them to try a service, not trust an app.", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# === SLIDE 12: Business Model ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "12 / BUSINESS MODEL", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Unit economics.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "Platform take-rate on every service.", size=18, color=GRAY)

# Unit economics table
metrics = [
    ("Avg service price", "$15"),
    ("Tech payout", "$10-11 (65-70%)"),
    ("Platform margin", "$4-5 (27-33%)"),
    ("Services/tech/day", "15-20"),
    ("Tech daily earnings", "$150-220")
]
y = 2.4
for metric, value in metrics:
    add_text_box(slide, 0.5, y, 3.5, 0.4, metric, size=13, color=WHITE)
    add_text_box(slide, 4.5, y, 2.5, 0.4, value, size=13, color=BLUE, bold=True)
    y += 0.45

# Margin stat box
add_stat_box(slide, 8, 2.4, 4, "27-33%", "Platform Margin", value_size=36)
add_text_box(slide, 8, 4.0, 4, 0.5, "High operating leverage — costs don't scale linearly", size=10, color=GRAY, align=PP_ALIGN.CENTER)

# At scale box
sbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.0), Inches(12), Inches(0.9))
sbox.fill.solid()
sbox.fill.fore_color.rgb = RgbColor(0, 40, 80)
sbox.line.color.rgb = BLUE
add_text_box(slide, 0.7, 5.2, 11.5, 0.5, "At Scale: 10,000 services/day = $150K revenue, $40-50K margin/day", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# === SLIDE 13: The Ask ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "13 / THE ASK", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Investment.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "$1.5M Seed to prove the model.", size=18, color=GRAY)

# Ask stat
add_stat_box(slide, 0.5, 2.4, 3.5, "$1.5M", "Seed Round", value_size=40)

# Use of funds
add_text_box(slide, 4.5, 2.4, 3, 0.4, "Use of Funds", size=14, color=BLUE, bold=True)
uses = [
    ("Product (40%)", "MVP app, API integrations"),
    ("Operations (30%)", "Techs, training, equipment"),
    ("Sales (20%)", "Enterprise BD"),
    ("G&A (10%)", "Legal, insurance")
]
y = 2.9
for label, desc in uses:
    add_text_box(slide, 4.5, y, 2.2, 0.35, label, size=11, color=WHITE, bold=True)
    add_text_box(slide, 6.8, y, 3, 0.35, desc, size=10, color=GRAY)
    y += 0.4

# Milestones
add_text_box(slide, 0.5, 4.5, 12, 0.4, "18-Month Milestones", size=14, color=BLUE, bold=True)
milestones = [
    ("First Paid Contract", "Month 6", "$10K+ MRR"),
    ("Product-Market Fit", "Month 12", "$100K MRR, 3 cities"),
    ("Series A Ready", "Month 18", "$300K MRR, 1,000+ vehicles")
]
y = 5.0
for milestone, timeline, target in milestones:
    add_text_box(slide, 0.5, y, 3.5, 0.35, milestone, size=12, color=WHITE)
    add_text_box(slide, 4.5, y, 2, 0.35, timeline, size=12, color=GRAY)
    add_text_box(slide, 7, y, 4, 0.35, target, size=12, color=GREEN, bold=True)
    y += 0.4

# === SLIDE 14: Exit ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 0.4, 4, 0.3, "14 / EXIT STRATEGY", size=10, color=BLUE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.8, "Returns.", size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.7, 12, 0.5, "Three paths to liquidity.", size=18, color=GRAY)

exits = [
    ("Acquisition (Most Likely)", "3-5 years, $100-500M", ["Waymo (uses 6+ contractors)", "Uber (just launched AV Services)", "Amazon (owns Zoox)", "", "Comp: Zoox = $1.2B"]),
    ("Fleet Services Roll-Up", "4-6 years", ["Enterprise Mobility", "Cox Automotive", "Fleetio ($1.5B valuation)"]),
    ("IPO Path", "7-10 years", ["Expand beyond AV to ALL", "commercial fleet services", "", "250K+ vehicles addressable"])
]
x = 0.5
for i, (title, timeline, points) in enumerate(exits):
    cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(3.9), Inches(2.8))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RgbColor(30, 30, 30)
    cbox.line.color.rgb = GREEN if i == 0 else RgbColor(60, 60, 60)
    add_text_box(slide, x + 0.15, 2.55, 3.6, 0.4, title, size=12, color=GREEN if i == 0 else WHITE, bold=True)
    add_text_box(slide, x + 0.15, 2.95, 3.6, 0.3, timeline, size=10, color=GRAY)
    y = 3.35
    for pt in points:
        if pt:
            add_text_box(slide, x + 0.15, y, 3.6, 0.3, "• " + pt if not pt.startswith("Comp") else pt, size=9, color=WHITE)
        y += 0.3
    x += 4.1

add_text_box(slide, 0.5, 5.5, 12, 0.4, "Waymo is valued at $126B. A $200M acquisition is a rounding error.", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# === SLIDE 15: Close ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

add_text_box(slide, 0.5, 2.2, 12.333, 0.8, "The future of mobility needs infrastructure.", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.2, 12.333, 0.5, "AUTONOMY, MAINTAINED.", size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 4.5, 12.333, 0.4, "Adonis Williams — Founder & CEO", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 5.0, 12.333, 0.35, "Deputy Functional Chief Engineer, Boeing (15 years)", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 5.4, 12.333, 0.35, "MS Systems Engineering | 2 Patents", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 6.0, 12.333, 0.35, "adonis@motra.io", size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Save
prs.save('/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Pitch-Deck-v3.pptx')
print("PowerPoint saved: MOTRA-Pitch-Deck-v3.pptx")
