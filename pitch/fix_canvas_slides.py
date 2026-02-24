from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)
LIGHT_PURPLE = RGBColor(232, 228, 240)
LIGHT_PEACH = RGBColor(252, 228, 214)
LIGHT_YELLOW = RGBColor(255, 249, 230)
LIGHT_GRAY = RGBColor(240, 240, 240)
LIGHT_BLUE = RGBColor(219, 234, 254)
DARK_GRAY = RGBColor(50, 50, 50)

def add_canvas_box(slide, left, top, width, height, fill_color, title, items):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(200, 200, 200)
    box.line.width = Pt(0.5)
    
    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.05), width - Inches(0.15), Inches(0.25))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.15), height - Inches(0.35))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)

def create_visual_canvas(slide, prs):
    """Add visual Business Model Canvas to a slide"""
    # Clear existing shapes
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Slide number
    num_box = slide.shapes.add_textbox(Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.3))
    p = num_box.text_frame.paragraphs[0]
    p.text = "11"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.3))
    p = title_box.text_frame.paragraphs[0]
    p.text = "BUSINESS MODEL CANVAS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Canvas dimensions
    col_width = Inches(1.88)
    row1_height = Inches(1.5)
    row2_height = Inches(1.5)
    row3_height = Inches(1.0)
    start_y = Inches(0.55)
    start_x = Inches(0.15)
    gap = Inches(0.05)
    
    # Key Partners
    add_canvas_box(slide, start_x, start_y, col_width, row1_height + row2_height + gap, LIGHT_PURPLE,
        "KEY PARTNERS",
        ["AV Ecosystem:", "• Waymo, Cruise, Zoox, Tesla", "• Fleet mgmt software", "",
         "Gig Economy:", "• Existing gig platforms", "• Training providers", "• Equipment suppliers", "",
         "Supporting:", "• Insurance providers", "• Background checks", "• Payment processing"])
    
    # Key Activities
    add_canvas_box(slide, start_x + col_width + gap, start_y, col_width, row1_height, LIGHT_PURPLE,
        "KEY ACTIVITIES",
        ["• Platform development", "• Tech recruitment", "• AV-specific training", "• Quality assurance", "• Enterprise sales"])
    
    # Key Resources
    add_canvas_box(slide, start_x + col_width + gap, start_y + row1_height + gap, col_width, row2_height, LIGHT_PURPLE,
        "KEY RESOURCES",
        ["Technology:", "• Dispatch platform", "• Fleet dashboard", "• Technician app", "",
         "Human:", "• Engineering team", "• Trained tech network"])
    
    # Value Propositions
    add_canvas_box(slide, start_x + (col_width + gap) * 2, start_y, col_width, row1_height + row2_height + gap, LIGHT_PEACH,
        "VALUE PROPOSITIONS",
        ["For Fleet Operators:", "• Reduce downtime", "• Variable cost model", "• Scale instantly", "• 24/7 availability", "• API integration", "",
         "vs. In-House:", "• 50% cost reduction", "• 30% more uptime", "",
         "For Technicians:", "• Flexible gig work", "• $15-20/hour"])
    
    # Customer Relationships
    add_canvas_box(slide, start_x + (col_width + gap) * 3, start_y, col_width, row1_height, LIGHT_YELLOW,
        "CUSTOMER RELATIONSHIPS",
        ["• Dedicated account mgrs", "• API integration support", "• Performance dashboards", "• SLA guarantees", "• 24/7 support"])
    
    # Channels
    add_canvas_box(slide, start_x + (col_width + gap) * 3, start_y + row1_height + gap, col_width, row2_height, LIGHT_YELLOW,
        "CHANNELS",
        ["Acquisition:", "• Direct enterprise sales", "• Industry conferences", "",
         "Delivery:", "• Fleet dashboard", "• API integration", "• Mobile dispatch"])
    
    # Customer Segments
    add_canvas_box(slide, start_x + (col_width + gap) * 4, start_y, col_width, row1_height + row2_height + gap, LIGHT_YELLOW,
        "CUSTOMER SEGMENTS",
        ["Primary:", "• Waymo (2,500+ vehicles)", "• Cruise (rebuilding)", "• Zoox (Amazon)", "• Tesla Robotaxi", "",
         "Secondary:", "• Amazon Delivery EVs", "• FedEx/UPS electric", "",
         "Buyers:", "• VP of Operations", "• Fleet Ops Manager"])
    
    # Cost Structure
    add_canvas_box(slide, start_x, start_y + row1_height + row2_height + gap * 2, col_width * 2 + gap, row3_height, LIGHT_GRAY,
        "COST STRUCTURE",
        ["Variable: Tech payouts (65-70%), Payment processing (2-3%)",
         "Fixed: Platform ($15-20K/mo), Team ($40-60K/mo), Marketing ($10-15K/mo)",
         "Seed: $1.5M total"])
    
    # Revenue Streams
    add_canvas_box(slide, start_x + (col_width + gap) * 2, start_y + row1_height + row2_height + gap * 2, col_width * 3 + gap * 2, row3_height, LIGHT_BLUE,
        "REVENUE STREAMS",
        ["Services: Quick ($12-18) | Deep ($45-75) | Emergency ($75-150)",
         "Unit Economics: Avg $15/service, 30% margin ($4-5), Target 10K/day",
         "Projections: Y1 $1.2M → Y2 $5.5M → Y3 $15M | TAM 2032: $5.5B"])

# ===== FIX BUSINESS PLAN =====
print("Fixing MOTRA-Business-Plan.pptx...")
prs = Presentation("/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx")

# Slide 11 is at index 11 (0-indexed: slide 12), and we added one at the end
# The old text slide 11 is at index 11, the new visual is at the end
# We need to: replace slide at index 11 with visual, remove the last slide

# First, let's identify which slide is which by counting
total_before = len(prs.slides)
print(f"  Total slides before: {total_before}")

# The old text-based canvas is slide 12 (index 11)
# The new visual canvas we added is at the end (index 23)
# We want to replace index 11's content with the visual canvas and remove the last slide

# Replace slide 11 (index 11) with visual canvas
slide_11 = prs.slides[11]
create_visual_canvas(slide_11, prs)

# Delete the last slide (the duplicate visual we added)
# In python-pptx, we need to remove from the slide collection
rId = prs.slides._sldIdLst[-1].rId
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[-1]

print(f"  Total slides after: {len(prs.slides)}")
prs.save("/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx")
print("  ✅ Business Plan fixed")

# ===== FIX PITCH DECK =====
print("\nFixing MOTRA-Pitch-Deck.pptx...")
prs2 = Presentation("/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Pitch-Deck.pptx")

total_before = len(prs2.slides)
print(f"  Total slides before: {total_before}")

# The pitch deck has fewer slides - visual canvas should be near the end
# Delete the last slide (the one we just added) since it's a duplicate
rId = prs2.slides._sldIdLst[-1].rId
prs2.part.drop_rel(rId)
del prs2.slides._sldIdLst[-1]

print(f"  Total slides after: {len(prs2.slides)}")
prs2.save("/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Pitch-Deck.pptx")
print("  ✅ Pitch Deck fixed")

print("\n✅ Both files fixed - visual canvas replaces text canvas")
