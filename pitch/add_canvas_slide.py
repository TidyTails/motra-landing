from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation("/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx")

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
BLUE = RGBColor(0, 102, 255)
WHITE = RGBColor(255, 255, 255)
LIGHT_PURPLE = RGBColor(232, 228, 240)
LIGHT_PEACH = RGBColor(252, 228, 214)
LIGHT_YELLOW = RGBColor(255, 249, 230)
LIGHT_GRAY = RGBColor(240, 240, 240)
LIGHT_BLUE = RGBColor(219, 234, 254)
DARK_GRAY = RGBColor(50, 50, 50)

def add_canvas_box(slide, left, top, width, height, fill_color, title, items):
    """Add a canvas section box"""
    # Box shape
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(200, 200, 200)
    box.line.width = Pt(0.5)
    
    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.05), width - Inches(0.15), Inches(0.25))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.15), height - Inches(0.35))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)

# Create new slide at position 12 (after North Star, before detailed breakdowns)
slide_layout = prs.slide_layouts[6]  # Blank layout

# Insert at position 11 (0-indexed, so this becomes slide 12)
# We'll add to end first, then note position
slide = prs.slides.add_slide(slide_layout)

# Header bar
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.5))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

# Slide number
num_box = slide.shapes.add_textbox(Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.3))
p = num_box.text_frame.paragraphs[0]
p.text = "11"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.3))
p = title_box.text_frame.paragraphs[0]
p.text = "BUSINESS MODEL CANVAS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE

# Canvas layout (5 columns, 3 rows)
# Row heights: top row = 1.6in, middle row = 1.6in, bottom row = 1.1in
# Total height: 4.3in (starts at 0.6in, ends at 4.9in)

col_width = Inches(1.88)  # ~9.4in / 5 columns
row1_height = Inches(1.5)
row2_height = Inches(1.5)
row3_height = Inches(1.0)
start_y = Inches(0.55)
start_x = Inches(0.15)
gap = Inches(0.05)

# ROW 1-2: Key Partners (spans 2 rows)
add_canvas_box(slide, start_x, start_y, col_width, row1_height + row2_height + gap, LIGHT_PURPLE,
    "KEY PARTNERS",
    ["AV Ecosystem:", "• Waymo, Cruise, Zoox, Tesla", "• Fleet mgmt software", "",
     "Gig Economy:", "• Existing gig platforms", "• Training providers", "• Equipment suppliers", "",
     "Supporting:", "• Insurance providers", "• Background checks", "• Payment processing"])

# ROW 1: Key Activities
add_canvas_box(slide, start_x + col_width + gap, start_y, col_width, row1_height, LIGHT_PURPLE,
    "KEY ACTIVITIES",
    ["• Platform development", "• Tech recruitment", "• AV-specific training", "• Quality assurance", "• Enterprise sales"])

# ROW 2: Key Resources
add_canvas_box(slide, start_x + col_width + gap, start_y + row1_height + gap, col_width, row2_height, LIGHT_PURPLE,
    "KEY RESOURCES",
    ["Technology:", "• Dispatch platform", "• Fleet dashboard", "• Technician app", "",
     "Human:", "• Engineering team", "• Trained tech network"])

# ROW 1-2: Value Propositions (spans 2 rows)
add_canvas_box(slide, start_x + (col_width + gap) * 2, start_y, col_width, row1_height + row2_height + gap, LIGHT_PEACH,
    "VALUE PROPOSITIONS",
    ["For Fleet Operators:", "• Reduce downtime", "• Variable cost model", "• Scale instantly", "• 24/7 availability", "• API integration", "",
     "vs. In-House:", "• 50% cost reduction", "• 30% more uptime", "",
     "For Technicians:", "• Flexible gig work", "• $15-20/hour"])

# ROW 1: Customer Relationships
add_canvas_box(slide, start_x + (col_width + gap) * 3, start_y, col_width, row1_height, LIGHT_YELLOW,
    "CUSTOMER RELATIONSHIPS",
    ["• Dedicated account mgrs", "• API integration support", "• Performance dashboards", "• SLA guarantees", "• 24/7 support"])

# ROW 2: Channels
add_canvas_box(slide, start_x + (col_width + gap) * 3, start_y + row1_height + gap, col_width, row2_height, LIGHT_YELLOW,
    "CHANNELS",
    ["Acquisition:", "• Direct enterprise sales", "• Industry conferences", "",
     "Delivery:", "• Fleet dashboard", "• API integration", "• Mobile dispatch"])

# ROW 1-2: Customer Segments (spans 2 rows)
add_canvas_box(slide, start_x + (col_width + gap) * 4, start_y, col_width, row1_height + row2_height + gap, LIGHT_YELLOW,
    "CUSTOMER SEGMENTS",
    ["Primary:", "• Waymo (2,500+ vehicles)", "• Cruise (rebuilding)", "• Zoox (Amazon)", "• Tesla Robotaxi", "",
     "Secondary:", "• Amazon Delivery EVs", "• FedEx/UPS electric", "",
     "Buyers:", "• VP of Operations", "• Fleet Ops Manager"])

# ROW 3: Cost Structure (spans 2 columns)
add_canvas_box(slide, start_x, start_y + row1_height + row2_height + gap * 2, col_width * 2 + gap, row3_height, LIGHT_GRAY,
    "COST STRUCTURE",
    ["Variable: Tech payouts (65-70%), Payment processing (2-3%)",
     "Fixed: Platform ($15-20K/mo), Team ($40-60K/mo), Marketing ($10-15K/mo)",
     "Seed: $1.5M total"])

# ROW 3: Revenue Streams (spans 3 columns)
add_canvas_box(slide, start_x + (col_width + gap) * 2, start_y + row1_height + row2_height + gap * 2, col_width * 3 + gap * 2, row3_height, LIGHT_BLUE,
    "REVENUE STREAMS",
    ["Services: Quick ($12-18) | Deep ($45-75) | Emergency ($75-150)",
     "Unit Economics: Avg $15/service, 30% margin ($4-5), Target 10K/day",
     "Projections: Y1 $1.2M → Y2 $5.5M → Y3 $15M | TAM 2032: $5.5B"])

# Save
output_path = "/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx"
prs.save(output_path)
print(f"✅ Added visual Business Model Canvas slide")
print(f"Total slides: {len(prs.slides)}")
