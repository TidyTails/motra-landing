#!/usr/bin/env python3
"""
MOTRA Pitch Deck v2 - PowerPoint Generator
Updated March 2026 with latest market data
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RgbColor(0, 102, 255)
DARK = RgbColor(26, 26, 26)
GRAY = RgbColor(102, 102, 102)
WHITE = RgbColor(255, 255, 255)
RED = RgbColor(220, 38, 38)
GREEN = RgbColor(22, 163, 74)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, slide_num, section, title, content_func):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK
    background.line.fill.background()
    
    # Slide number
    num_box = slide.shapes.add_textbox(Inches(12), Inches(0.3), Inches(1), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num:02d} / {section.upper()}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT
    
    # Section label
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.3))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = section
    p.font.size = Pt(11)
    p.font.color.rgb = BLUE
    p.font.bold = True
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    content_func(slide)
    
    return slide

def add_text_box(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box

def add_stat_box(slide, left, top, value, label):
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.5), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(0, 102, 255)
    box.fill.fore_color.brightness = 0.85
    box.line.color.rgb = BLUE
    
    # Value
    add_text_box(slide, left, top + 0.1, 2.5, 0.6, value, size=36, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left, top + 0.7, 2.5, 0.4, label, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# === SLIDES ===

# Slide 0: Title
add_title_slide(prs, "MOTRA", "AUTONOMY, MAINTAINED")

# Slide 1: Company Purpose
def content_purpose(slide):
    add_text_box(slide, 0.5, 2.2, 12, 1.5, 
        "MOTRA is the infrastructure layer for autonomous vehicle fleet care.",
        size=28, color=WHITE)
    add_text_box(slide, 0.5, 3.5, 12, 1,
        "We deploy a gig-powered network of mobile technicians to clean and service robotaxis on-location, on-demand — keeping fleets running 24/7.",
        size=18, color=GRAY)
    # Highlight box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(12), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.fill.fore_color.brightness = 0.85
    box.line.color.rgb = BLUE
    add_text_box(slide, 0.7, 5.2, 11.5, 0.6, "Think DoorDash for fleet maintenance — invisible, essential, everywhere.", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_content_slide(prs, 1, "Company Purpose", "The infrastructure layer for AV fleet care.", content_purpose)

# Slide 2: Problem
def content_problem(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.8,
        "Robotaxis run 15+ hours a day. Their maintenance infrastructure doesn't.",
        size=24, color=WHITE)
    
    # Pain points
    points = [
        ("No driver = no eyes", "Nobody notices trash, spills, or wear between rides"),
        ("Depot-dependent", "Vehicles must return to facilities for basic cleaning"),
        ("Downtime = lost revenue", "Every minute in depot is a ride not taken"),
        ("Sensors are critical", "Dirty LiDAR/cameras = degraded driving")
    ]
    
    y = 3.2
    for i, (title, desc) in enumerate(points):
        add_text_box(slide, 0.5, y, 0.4, 0.4, str(i+1), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, 1.0, y, 5, 0.3, title, size=16, color=WHITE, bold=True)
        add_text_box(slide, 1.0, y + 0.35, 5, 0.3, desc, size=12, color=GRAY)
        y += 0.9
    
    # Current reality box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3.2), Inches(5.5), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(40, 40, 40)
    box.line.fill.background()
    add_text_box(slide, 7.2, 3.4, 5, 0.4, "Current Reality", size=16, color=RED, bold=True)
    add_text_box(slide, 7.2, 3.9, 5, 1.5, "Waymo uses 1,000+ Transdev operators at centralized depots.\n\nThat doesn't scale to 1 million rides/week.", size=14, color=WHITE)

add_content_slide(prs, 2, "The Problem", "The maintenance bottleneck.", content_problem)

# Slide 3: Solution
def content_solution(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6,
        "Mobile fleet care, dispatched like an Uber.",
        size=28, color=WHITE)
    add_text_box(slide, 0.5, 2.9, 12, 0.5,
        "MOTRA deploys certified technicians directly to vehicles — wherever they are.",
        size=16, color=GRAY)
    
    # Service boxes
    services = [
        ("Quick Clean", "$12-18", "5-10 min wipe-down, trash, odor"),
        ("Deep Clean", "$45-75", "Full interior detail, sensors, exterior"),
        ("Emergency", "$75-150", "On-demand response for incidents")
    ]
    
    x = 0.5
    for title, price, desc in services:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.8), Inches(3.8), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(30, 30, 30)
        box.line.color.rgb = RgbColor(50, 50, 50)
        add_text_box(slide, x + 0.2, 4.0, 3.4, 0.4, title, size=18, color=WHITE, bold=True)
        add_text_box(slide, x + 0.2, 4.5, 3.4, 0.4, price, size=24, color=BLUE, bold=True)
        add_text_box(slide, x + 0.2, 5.1, 3.4, 0.5, desc, size=12, color=GRAY)
        x += 4.1

add_content_slide(prs, 3, "The Solution", "Mobile fleet care on demand.", content_solution)

# Slide 4: Why Now
def content_why_now(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6,
        "The robotaxi industry just hit escape velocity.",
        size=28, color=WHITE)
    
    # Stats row
    stats = [
        ("400K", "Rides/Week"),
        ("10", "US Markets"),
        ("$16B", "Fresh Funding"),
        ("$126B", "Valuation")
    ]
    x = 0.5
    for value, label in stats:
        add_stat_box(slide, x, 3.0, value, label)
        x += 3.0
    
    # Table
    add_text_box(slide, 0.5, 4.5, 2, 0.3, "Metric", size=11, color=BLUE, bold=True)
    add_text_box(slide, 3, 4.5, 3, 0.3, "Now (Feb 2026)", size=11, color=BLUE, bold=True)
    add_text_box(slide, 7, 4.5, 4, 0.3, "EOY 2026 Target", size=11, color=BLUE, bold=True)
    
    rows = [
        ("Weekly Rides", "400,000", "1,000,000"),
        ("Markets", "10 cities", "20+ cities (incl. London, Tokyo)"),
        ("Fleet Size", "2,500 vehicles", "4,500+ vehicles"),
        ("Price War", "Waymo $8.17 vs Uber $17.25", "Volume exploding")
    ]
    y = 4.9
    for metric, now, target in rows:
        add_text_box(slide, 0.5, y, 2, 0.35, metric, size=12, color=WHITE)
        add_text_box(slide, 3, y, 3.5, 0.35, now, size=12, color=WHITE, bold=True)
        add_text_box(slide, 7, y, 5, 0.35, target, size=12, color=WHITE)
        y += 0.45
    
    add_text_box(slide, 0.5, 6.8, 12, 0.3, "Sources: TechXplore Feb 24 2026, Fox News Feb 27 2026", size=9, color=GRAY)

add_content_slide(prs, 4, "Why Now", "Escape velocity.", content_why_now)

# Slide 5: Market Size
def content_market(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "Robotaxi market: 99% CAGR through 2033", size=28, color=WHITE)
    
    # Big stats
    stats = [("$2.7B", "2025 TAM"), ("$147B", "2033 TAM"), ("99%", "CAGR")]
    x = 0.5
    for value, label in stats:
        add_stat_box(slide, x, 3.0, value, label)
        x += 4.0
    
    # MOTRA TAM
    add_text_box(slide, 0.5, 4.6, 12, 0.4, "MOTRA's Serviceable Market", size=18, color=BLUE, bold=True)
    rows = [
        ("2028", "$547M", "50K robotaxis × $30/week"),
        ("2030", "$2.2B", "200K robotaxis"),
        ("2032", "$5.5B", "500K robotaxis + commercial EVs")
    ]
    y = 5.1
    for year, tam, note in rows:
        add_text_box(slide, 0.5, y, 1.5, 0.35, year, size=14, color=WHITE)
        add_text_box(slide, 2, y, 2, 0.35, tam, size=14, color=BLUE, bold=True)
        add_text_box(slide, 4.5, y, 7, 0.35, note, size=12, color=GRAY)
        y += 0.45

add_content_slide(prs, 5, "Market Size", "A $147B opportunity.", content_market)

# Slide 6: Competition
def content_competition(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "Competitors are emerging. We have the edge.", size=28, color=WHITE)
    
    # Competition table
    add_text_box(slide, 0.5, 3.0, 3, 0.3, "Player", size=11, color=BLUE, bold=True)
    add_text_box(slide, 4, 3.0, 4, 0.3, "Approach", size=11, color=BLUE, bold=True)
    add_text_box(slide, 8.5, 3.0, 4, 0.3, "Limitation", size=11, color=BLUE, bold=True)
    
    rows = [
        ("Uber AV Services", "Depot-based cleaning", "Competitor to Waymo"),
        ("Transdev", "1,000+ operators at depots", "Depot-bound, not mobile"),
        ("Tesla Robots", "Automated cleaning", "Years away, Tesla-only"),
        ("In-House Ops", "AV companies DIY", "High fixed cost")
    ]
    y = 3.5
    for player, approach, limitation in rows:
        add_text_box(slide, 0.5, y, 3.5, 0.4, player, size=13, color=WHITE, bold=True)
        add_text_box(slide, 4, y, 4, 0.4, approach, size=12, color=GRAY)
        add_text_box(slide, 8.5, y, 4, 0.4, limitation, size=12, color=RED)
        y += 0.55
    
    # MOTRA advantage box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = GREEN
    box.fill.fore_color.brightness = 0.85
    box.line.color.rgb = GREEN
    add_text_box(slide, 0.7, 6.0, 11.5, 0.6, "MOTRA Advantage: Neutral platform (Switzerland) · Mobile-first · 100% focused on AV fleet care", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_content_slide(prs, 6, "Competition", "Our competitive edge.", content_competition)

# Slide 7: Why Won't Waymo Build This
def content_not_inhouse(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "Why won't Waymo just do this themselves?", size=28, color=WHITE)
    
    reasons = [
        ("1. They Already Outsource", "Waymo uses Transdev, Moove, Avis, TechForce, Amerit, Terawatt today. Build vs. buy = they buy."),
        ("2. Not Core Competency", "Waymo's edge is autonomous driving AI. Every $ on cleaning = $ not on AI."),
        ("3. Math Doesn't Work", "20+ cities = thousands of employees, facilities, HR. Gig model = variable cost.")
    ]
    
    x = 0.5
    for title, desc in reasons:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.2), Inches(3.8), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(30, 30, 30)
        box.line.color.rgb = BLUE
        add_text_box(slide, x + 0.2, 3.4, 3.4, 0.4, title, size=14, color=BLUE, bold=True)
        add_text_box(slide, x + 0.2, 3.9, 3.4, 1.4, desc, size=12, color=WHITE)
        x += 4.1
    
    # Highlight
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.fill.fore_color.brightness = 0.85
    box.line.color.rgb = BLUE
    add_text_box(slide, 0.7, 6.0, 11.5, 0.6, "Uber just validated this. If big players could do it better in-house, they wouldn't be building outsourced solutions.", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_content_slide(prs, 7, "Why Not In-House", "The outsourcing thesis.", content_not_inhouse)

# Slide 8: Our Moat
def content_moat(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "Three layers of defensibility.", size=28, color=WHITE)
    
    moats = [
        ("Year 1-2: Network Effects", "• First AV-specialized gig network per city\n• Hyperlocal network effects\n• More techs → faster response → more volume"),
        ("Year 2-3: Data Advantage", "• Every service = data on times, issues, vehicles\n• Predictive maintenance insights\n• DoorDash has 30% of eng in data roles"),
        ("Year 3+: Integration Lock-In", "• API embedded in fleet management\n• Part of operator workflows\n• 5% retention = 25-95% profit increase")
    ]
    
    x = 0.5
    for title, points in moats:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.2), Inches(3.8), Inches(2.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(30, 30, 30)
        box.line.color.rgb = BLUE
        add_text_box(slide, x + 0.2, 3.4, 3.4, 0.5, title, size=13, color=BLUE, bold=True)
        add_text_box(slide, x + 0.2, 4.0, 3.4, 1.8, points, size=11, color=WHITE)
        x += 4.1
    
    add_text_box(slide, 0.5, 6.3, 12, 0.5, "The Uber analogy: Uber's moat isn't the app — it's the network of drivers in every city.", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_content_slide(prs, 8, "Defensibility", "Our moat.", content_moat)

# Slide 9: Team
def content_team(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "15 years at Boeing. 2 patents. Built to scale.", size=28, color=WHITE)
    
    # Team box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.2), Inches(5.5), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(30, 30, 30)
    box.line.color.rgb = BLUE
    
    add_text_box(slide, 0.7, 3.4, 5, 0.5, "Adonis Williams — Founder & CEO", size=18, color=WHITE, bold=True)
    
    details = [
        ("Current Role", "Deputy Functional Chief Engineer, Boeing"),
        ("Tenure", "14 years 8 months at Boeing"),
        ("Patents", "2 patents (cybersecurity, blockchain comms)"),
        ("Education", "MS Systems Engineering, Missouri S&T"),
        ("Experience", "Distributed teams, USAF pilot tests, cyber")
    ]
    y = 4.0
    for label, value in details:
        add_text_box(slide, 0.7, y, 1.8, 0.35, label, size=10, color=GRAY)
        add_text_box(slide, 2.6, y, 3, 0.35, value, size=11, color=WHITE)
        y += 0.4
    
    # Why Adonis box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3.2), Inches(6), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(20, 20, 30)
    box.line.fill.background()
    
    add_text_box(slide, 6.7, 3.4, 5.5, 0.4, "Why Adonis Wins This", size=16, color=BLUE, bold=True)
    add_text_box(slide, 6.7, 3.9, 5.5, 1.5, 
        "MOTRA is a systems problem: coordinating a distributed workforce to service safety-critical autonomous vehicles.\n\nThat's exactly what he's been doing at Boeing for 15 years — applied to mobility instead of aviation.",
        size=12, color=WHITE)
    add_text_box(slide, 6.7, 5.6, 5.5, 1,
        "✓ Complex operations at scale\n✓ Distributed teams across geographies\n✓ Safety-critical systems\n✓ Innovation track record (2 patents)",
        size=11, color=GREEN)

add_content_slide(prs, 9, "Team", "Leadership.", content_team)

# Slide 10: Risk Mitigation
def content_risk(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "What if the AV market slows down?", size=28, color=WHITE)
    
    # Risk box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.2), Inches(5.5), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(50, 30, 30)
    box.line.color.rgb = RED
    add_text_box(slide, 0.7, 3.4, 5, 0.4, "The Risk", size=14, color=RED, bold=True)
    add_text_box(slide, 0.7, 3.9, 5, 0.7, "Robotaxi adoption could hit regulatory delays, technical hurdles, or economic headwinds.", size=12, color=WHITE)
    
    # Mitigation box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3.2), Inches(6), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(30, 50, 30)
    box.line.color.rgb = GREEN
    add_text_box(slide, 6.7, 3.4, 5.5, 0.4, "Adjacent Markets", size=14, color=GREEN, bold=True)
    add_text_box(slide, 6.7, 3.9, 5.5, 0.7, "Even if AV slows, 250K+ commercial EVs need regular cleaning and maintenance.", size=12, color=WHITE)
    
    # Fleet table
    add_text_box(slide, 0.5, 5.0, 3, 0.3, "Fleet", size=11, color=BLUE, bold=True)
    add_text_box(slide, 4, 5.0, 3, 0.3, "Current", size=11, color=BLUE, bold=True)
    add_text_box(slide, 7.5, 5.0, 4, 0.3, "Target", size=11, color=BLUE, bold=True)
    
    fleets = [
        ("Amazon (Rivian)", "30,000 EVs", "100,000 by 2030"),
        ("FedEx", "200,000+ vehicles", "All-electric by 2040"),
        ("UPS", "18,300+ alt-fuel", "10,000 more ordered")
    ]
    y = 5.4
    for fleet, current, target in fleets:
        add_text_box(slide, 0.5, y, 3.5, 0.35, fleet, size=13, color=WHITE, bold=True)
        add_text_box(slide, 4, y, 3, 0.35, current, size=12, color=WHITE)
        add_text_box(slide, 7.5, y, 4, 0.35, target, size=12, color=WHITE)
        y += 0.45
    
    add_text_box(slide, 0.5, 6.6, 12, 0.4, "MOTRA has multiple paths to scale. Robotaxis are the beachhead — commercial EVs are the expansion.", size=14, color=GREEN, align=PP_ALIGN.CENTER)

add_content_slide(prs, 10, "Risk Mitigation", "De-risked.", content_risk)

# Slide 11: GTM
def content_gtm(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "The DoorDash playbook: do things that don't scale.", size=28, color=WHITE)
    
    phases = [
        ("Phase 1: Manual Pilot", "Month 1-3", "• ONE operator, ONE city\n• Free 30-day pilot, 50 vehicles\n• Text dispatch, Google Sheets\n• We ARE the operators"),
        ("Phase 2: Paid Contract", "Month 4-6", "• Convert to paid\n• Price 20% below Transdev\n• 6-month commitment\n• Build reference customer"),
        ("Phase 3: Scale", "Month 7-12", "• Add second city\n• Launch MVP app\n• Use first customer as ref\n• Target $100K MRR")
    ]
    
    x = 0.5
    for title, timeline, points in phases:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.2), Inches(3.8), Inches(2.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(30, 30, 30)
        box.line.color.rgb = BLUE
        add_text_box(slide, x + 0.2, 3.4, 3.4, 0.4, title, size=14, color=BLUE, bold=True)
        add_text_box(slide, x + 0.2, 3.85, 3.4, 0.3, timeline, size=11, color=GRAY)
        add_text_box(slide, x + 0.2, 4.2, 3.4, 1.6, points, size=11, color=WHITE)
        x += 4.1
    
    add_text_box(slide, 0.5, 6.3, 12, 0.5, "Warm intros = 3-5x better close rate. We're not asking them to trust an app — we're asking them to try a service.", size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_content_slide(prs, 11, "Go-To-Market", "First customers.", content_gtm)

# Slide 12: Business Model
def content_biz_model(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "Platform take-rate on every service.", size=28, color=WHITE)
    
    # Unit economics table
    metrics = [
        ("Avg service price", "$15"),
        ("Tech payout", "$10-11 (65-70%)"),
        ("Platform margin", "$4-5 (27-33%)"),
        ("Services/tech/day", "15-20"),
        ("Tech daily earnings", "$150-220")
    ]
    y = 3.2
    for metric, value in metrics:
        add_text_box(slide, 0.5, y, 3.5, 0.4, metric, size=14, color=WHITE)
        add_text_box(slide, 4.5, y, 2.5, 0.4, value, size=14, color=BLUE, bold=True)
        y += 0.5
    
    # Margin stat
    add_stat_box(slide, 8, 3.2, "25-35%", "Platform Margin")
    add_text_box(slide, 8, 4.6, 4, 0.6, "High operating leverage — platform costs don't scale linearly", size=11, color=GRAY, align=PP_ALIGN.CENTER)
    
    # At scale box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.fill.fore_color.brightness = 0.85
    box.line.color.rgb = BLUE
    add_text_box(slide, 0.7, 6.0, 11.5, 0.6, "At Scale: 10,000 services/day = $150K revenue, $40-50K margin/day", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_content_slide(prs, 12, "Business Model", "Unit economics.", content_biz_model)

# Slide 13: The Ask
def content_ask(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "$1.5M Seed to prove the model.", size=28, color=WHITE)
    
    # Big ask stat
    add_stat_box(slide, 0.5, 3.2, "$1.5M", "Seed Round")
    
    # Use of funds
    add_text_box(slide, 3.5, 3.2, 3, 0.4, "Use of Funds", size=14, color=BLUE, bold=True)
    uses = [
        ("Product (40%)", "MVP app, API integrations"),
        ("Operations (30%)", "Techs, training, equipment"),
        ("Sales (20%)", "Enterprise BD"),
        ("G&A (10%)", "Legal, insurance")
    ]
    y = 3.7
    for label, desc in uses:
        add_text_box(slide, 3.5, y, 2, 0.35, label, size=11, color=WHITE, bold=True)
        add_text_box(slide, 5.5, y, 3, 0.35, desc, size=11, color=GRAY)
        y += 0.4
    
    # Milestones
    add_text_box(slide, 0.5, 5.4, 12, 0.4, "18-Month Milestones", size=16, color=BLUE, bold=True)
    milestones = [
        ("First Paid Contract", "Month 6", "$10K+ MRR"),
        ("Product-Market Fit", "Month 12", "$100K MRR, 3 cities"),
        ("Series A Ready", "Month 18", "$300K MRR, 1,000+ vehicles")
    ]
    y = 5.9
    for milestone, timeline, target in milestones:
        add_text_box(slide, 0.5, y, 3.5, 0.35, milestone, size=13, color=WHITE)
        add_text_box(slide, 4.5, y, 2, 0.35, timeline, size=13, color=GRAY)
        add_text_box(slide, 7, y, 4, 0.35, target, size=13, color=GREEN, bold=True)
        y += 0.45

add_content_slide(prs, 13, "The Ask", "Investment.", content_ask)

# Slide 14: Exit
def content_exit(slide):
    add_text_box(slide, 0.5, 2.2, 12, 0.6, "Three paths to liquidity.", size=28, color=WHITE)
    
    exits = [
        ("Acquisition (Most Likely)", "3-5 years, $100-500M", "• Waymo (uses 6+ contractors)\n• Uber (just launched AV Services)\n• Amazon (owns Zoox)\n\nComp: Zoox = $1.2B"),
        ("Fleet Services Roll-Up", "4-6 years", "• Enterprise Mobility\n• Cox Automotive\n• Fleetio ($1.5B valuation)"),
        ("IPO Path", "7-10 years", "Expand beyond AV to ALL commercial fleet services.\n\n250K+ vehicles addressable today.")
    ]
    
    x = 0.5
    colors = [GREEN, WHITE, WHITE]
    for i, (title, timeline, points) in enumerate(exits):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.2), Inches(3.8), Inches(3))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(30, 30, 30)
        box.line.color.rgb = GREEN if i == 0 else RgbColor(50, 50, 50)
        add_text_box(slide, x + 0.2, 3.4, 3.4, 0.5, title, size=13, color=colors[i], bold=True)
        add_text_box(slide, x + 0.2, 3.9, 3.4, 0.35, timeline, size=11, color=GRAY)
        add_text_box(slide, x + 0.2, 4.3, 3.4, 1.7, points, size=10, color=WHITE)
        x += 4.1
    
    add_text_box(slide, 0.5, 6.5, 12, 0.5, "Waymo is valued at $126B. A $200M acquisition is a rounding error.", size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_content_slide(prs, 14, "Exit Strategy", "Returns.", content_exit)

# Slide 15: Close
def content_close(slide):
    add_text_box(slide, 0.5, 2.8, 12, 0.8, "The future of mobility needs infrastructure.", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 3.8, 12, 0.5, "AUTONOMY, MAINTAINED.", size=20, color=GRAY, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.5, 5.2, 12, 0.4, "Adonis Williams — Founder & CEO", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 5.7, 12, 0.35, "Deputy Functional Chief Engineer, Boeing (14 years)", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 6.1, 12, 0.35, "MS Systems Engineering | 2 Patents", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 6.6, 12, 0.35, "adonis@motra.io", size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

add_content_slide(prs, 15, "", "Thank you.", content_close)

# Save
prs.save('/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Pitch-Deck-v2.pptx')
print("PowerPoint saved: MOTRA-Pitch-Deck-v2.pptx")
