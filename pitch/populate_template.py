from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Load the template
prs = Presentation("/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx")

# MOTRA Content
MOTRA_CONTENT = {
    # Slide 1 - Title
    "UW 540 / 440 Business Plan practicum": "MOTRA Business Plan",
    "Solving real problems and having a positive impact on the worl": "Autonomy, Maintained",
    
    # Slide 2 - Team
    "Insert Name": "Adonis Williams",
    "Elvis Chan\nEnweier Chau\t\nMinjay Jo": "Adonis Williams\nFounder & CEO",
    "Team player\nLeader\nComputer engineer  ": "Deputy Functional Chief Engineer, Boeing\nEngineering Manager, Boeing Research & Technology\nMissouri University of Science and Technology",
    "Role / Function on the Team": "Founder & CEO",
    "Brief Bio": "10+ years engineering leadership at Boeing, expertise in complex systems and fleet operations",
    "Identify a role for each team member": "Seattle, WA - Epicenter of tech and mobility innovation",
    
    # Slide 3 - Values
    "Accountability\nTransparency\nCollaboration\nUser-Centricity\nAgility": "Engineering Excellence\nScalability First\nOperator Mindset\nData-Driven\nCustomer Obsession",
    "Own Your Role\nTimely Updates\nConstructive Feedback\nFollow Through\nRetrospectives": "Own outcomes end-to-end\nDaily standups\nTest assumptions rapidly\nMeasure everything\nIterate weekly",
    "Data-Informed\nClarity on Authority\nDisagree & Commit\nDocument Decisions\nTime-Box Discussions": "Evidence over opinion\nFounder makes final calls\nMove fast once decided\nLog all pivots\n24-hour decision SLA",
    "Share work early and often in appropriate channels.\nSchedule meetings within core overlapping hours.": "Slack for async\nWeekly all-hands\nBi-weekly investor updates\nMonthly board reviews",
    
    # Slide 4 - Define the need
    "185 Seattle bars contacted expressing interest, with complete data collected (TV counts, channels, s": "2,500+ Waymo robotaxis operating 24/7, requiring constant cleaning and maintenance between rides. Current depot-based operations don't scale.",
    "Search any game (NFL, NBA, UFC, soccer, esports, WNBA, college) and instantly see every bar streamin": "MOTRA deploys gig-powered mobile technicians to clean and service autonomous vehicles on-location, on-demand — eliminating depot downtime.",
    
    # Slide 7 - Facts & Assumptions
    "2025: 5 paying bars ($60/mo) = $4K.": "Waymo does 400,000+ rides/week, targeting 1M/week by end of 2026",
    "2026: World Cup drives 100 bars + 20% boost adoption = $86K.": "Fleet growing from 2,500 to \"tens of thousands\" of vehicles",
    "2027-29: Expand to Portland, Vancouver BC, San Francisco with 30% subscription adoption.": "Tesla Robotaxi, Zoox, Cruise all scaling 2026-2027",
    "Costs: Hosting $50/mo, API $200/mo, marketing/tools $300/mo, team compensation scales with revenue.": "Service cost: $15/service, 27-33% platform margin, $4-5 per service profit",
    "Users (sports fans, free), Customers (sports bars seeking traffic), Buyers (bar owners/managers maki": "Users: AV Fleet Operators (Waymo, Cruise, Zoox)\nCustomers: Fleet Operations Managers\nBuyers: VP of Operations / Head of Fleet",
    "TAM: 2 million potential adult users in greater Seattle area | SAM: 240 sports bars in Seattle | SOM": "TAM: $5.5B by 2032 (500K AVs × 2 services/day × $15)\nSAM: $547M by 2028 (50K AVs)\nSOM: $2M Year 1 (10% of one operator in one city)",
    
    # Slide 16 - North Star
    "In the hustle-and-bustle and the ups and downs of your entrepreneurial journey, it is very helpful t": "Every autonomous vehicle in the world is serviced by MOTRA's invisible infrastructure layer — clean, maintained, and always ready.",
    
    # Business Canvas sections
    "TBD": "See detailed sections below",
}

# Function to replace text in shapes
def replace_text_in_shape(shape, old_text, new_text):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

# Go through all slides and replace text
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            for old_text, new_text in MOTRA_CONTENT.items():
                if old_text[:30] in full_text[:50]:  # Match beginning
                    # Replace text in all paragraphs/runs
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old, new in MOTRA_CONTENT.items():
                                if old in run.text:
                                    run.text = run.text.replace(old, new)

# Save
output_path = "/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Business-Plan.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
