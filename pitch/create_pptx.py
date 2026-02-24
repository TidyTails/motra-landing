from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RGBColor(0, 102, 255)
DARK_BG = RGBColor(10, 10, 10)
GRAY = RGBColor(136, 136, 136)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(102, 102, 102)

def add_slide(title_text, subtitle_text=None, content=None, is_title_slide=False):
    """Add a slide with consistent styling"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    if is_title_slide:
        # Title slide - centered
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.5))
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle_text
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = LIGHT_GRAY
            sub_para.alignment = PP_ALIGN.CENTER
    else:
        # Section label
        if subtitle_text:
            label_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(3), Inches(0.3))
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = subtitle_text.upper()
            label_para.font.size = Pt(11)
            label_para.font.color.rgb = BLUE
            label_para.font.bold = True
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.9), Inches(11.5), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Content area
        if content:
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.3), Inches(11.5), Inches(4.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, line in enumerate(content):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                para.text = line
                para.font.size = Pt(18)
                para.font.color.rgb = GRAY
                para.space_after = Pt(12)
    
    return slide

def add_stat_slide(title_text, subtitle_text, stats):
    """Add a slide with statistics"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Section label
    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(3), Inches(0.3))
    label_para = label_box.text_frame.paragraphs[0]
    label_para.text = subtitle_text.upper()
    label_para.font.size = Pt(11)
    label_para.font.color.rgb = BLUE
    label_para.font.bold = True
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.9), Inches(11.5), Inches(1))
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Stats grid
    start_x = 0.75
    stat_width = 2.8
    gap = 0.3
    
    for i, (value, label) in enumerate(stats):
        x = start_x + i * (stat_width + gap)
        
        # Stat box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(stat_width), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(17, 17, 17)
        box.line.color.rgb = RGBColor(34, 34, 34)
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(2.7), Inches(stat_width), Inches(1))
        val_para = val_box.text_frame.paragraphs[0]
        val_para.text = value
        val_para.font.size = Pt(44)
        val_para.font.bold = True
        val_para.font.color.rgb = BLUE
        val_para.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(stat_width), Inches(0.5))
        lbl_para = lbl_box.text_frame.paragraphs[0]
        lbl_para.text = label
        lbl_para.font.size = Pt(14)
        lbl_para.font.color.rgb = LIGHT_GRAY
        lbl_para.alignment = PP_ALIGN.CENTER
    
    return slide

# SLIDE 0: Title
add_slide("MOTRA", "Autonomy, Maintained", is_title_slide=True)

# SLIDE 1: Company Purpose
add_slide(
    "MOTRA is the infrastructure layer for autonomous vehicle fleet care.",
    "Company Purpose",
    [
        "We deploy a gig-powered network of mobile technicians to clean and service robotaxis on-location, on-demand.",
        "",
        "Think AWS for autonomous fleet maintenance — invisible, essential, everywhere."
    ]
)

# SLIDE 2: The Problem
add_slide(
    "Autonomous vehicles run 24/7. Their maintenance infrastructure doesn't.",
    "The Problem",
    [
        "1. No driver = no eyes — Nobody notices trash, spills, or wear between rides",
        "",
        "2. Depot-dependent — Vehicles must return to facilities for basic cleaning", 
        "",
        "3. Downtime = lost revenue — Every minute in depot is a ride not taken",
        "",
        "4. Sensors are safety-critical — Dirty LiDAR/cameras = degraded driving",
        "",
        "Current reality: Waymo hand-washes every vehicle at centralized depots.",
        "That doesn't scale to millions of rides."
    ]
)

# SLIDE 3: The Solution
add_slide(
    "Mobile fleet care, dispatched like an Uber.",
    "The Solution",
    [
        "MOTRA deploys certified technicians directly to vehicles — wherever they are.",
        "",
        "• Quick Clean (5-10 min) — Between-ride wipe-down, trash removal, odor neutralization",
        "",
        "• Deep Clean (30-60 min) — Full interior detail, sensor cleaning, exterior wash",
        "",
        "• Maintenance — Light repairs, tire checks, fluid top-offs, emergency response",
        "",
        "Result: Vehicles stay in service zones | Variable cost model | Scales instantly | 24/7 availability"
    ]
)

# SLIDE 4: Why Now
add_slide(
    "The robotaxi industry is hitting an inflection point.",
    "Why Now",
    [
        "Waymo Scale: 400K rides/week — targeting 1M/week by end of 2026",
        "",
        "Fleet Growth: 2,500 vehicles today — \"tens of thousands\" planned",
        "",
        "New Entrants: Tesla Robotaxi, Zoox, Cruise rebuild — all scaling 2026-2027",
        "",
        "Gig Infrastructure: Uber/DoorDash proved the model — workforce is trained and ready",
        "",
        "The window is NOW — before AV companies build in-house or a competitor emerges."
    ]
)

# SLIDE 5: Market Size
add_stat_slide(
    "$500M+ market emerging, scaling to $5B+",
    "Market Size",
    [
        ("$547M", "TAM 2028"),
        ("$2.2B", "TAM 2030"),
        ("$5.5B", "TAM 2032"),
        ("$2M", "Year 1 SOM")
    ]
)

# SLIDE 6: Competition
add_slide(
    "Blue ocean with fragmented alternatives.",
    "Competition",
    [
        "AV In-House Ops → High fixed cost, doesn't scale",
        "",
        "Traditional Fleet Services → Not mobile, not AV-specialized",
        "",
        "Car Washes → Damages sensors, no interior service",
        "",
        "Mobile Detailing → Not scaled, not fleet-focused",
        "",
        "Our Advantages:",
        "• AV-Specialized (sensors, EVs, safety)",
        "• Mobile-First (go to the vehicle)",
        "• Gig-Powered (variable cost, instant scale)",
        "• First Mover (no scaled competitor exists)"
    ]
)

# SLIDE 7: Product
add_slide(
    "Platform + Network + Expertise",
    "Product",
    [
        "For Fleet Operators:",
        "• Fleet Dashboard — real-time status, scheduling, quality metrics",
        "• API Integration — connects to existing fleet management systems",
        "• Analytics — predictive maintenance, cost tracking",
        "",
        "For MOTRA Techs:",
        "• Mobile App — job dispatch, checklists, earnings",
        "• Certification — AV-specific training program",
        "• Equipment Kits — standardized tools",
        "",
        "Roadmap: MVP Q3 2026 → Waymo API Q4 2026 → 3 Markets Q1 2027"
    ]
)

# SLIDE 8: Business Model
add_slide(
    "Platform take-rate on every service.",
    "Business Model",
    [
        "Average service price: $15",
        "Tech payout: $10-11",
        "Platform margin: $4-5 (27-33%)",
        "",
        "Services per tech per day: 15-20",
        "Tech daily earnings: $150-220",
        "",
        "At Scale: 10,000 services/day = $150K revenue, $40-50K margin",
        "",
        "High operating leverage — platform costs don't scale linearly with services."
    ]
)

# SLIDE 9: Team
add_slide(
    "Built to win this market.",
    "Team",
    [
        "Adonis Williams — Founder & CEO",
        "",
        "• Deputy Functional Chief Engineer at Boeing",
        "• Engineering Manager, Boeing Research & Technology",
        "• Deep expertise in complex systems, fleet operations, and scaling infrastructure",
        "• Education: Missouri University of Science and Technology",
        "• Based in Seattle, WA — epicenter of tech and mobility innovation",
        "",
        "Why We Win: Engineering Excellence • Systems Thinking • 100% Focus on AV Fleet Care"
    ]
)

# SLIDE 10: The Ask
add_slide(
    "Path to $10M ARR in 36 months.",
    "Financials & The Ask",
    [
        "Year 1: 1 market, 500 vehicles, $1.2M revenue",
        "Year 2: 3 markets, 3,000 vehicles, $5.5M revenue", 
        "Year 3: 7 markets, 10,000 vehicles, $15M revenue",
        "",
        "Seed Round: $1.5M",
        "",
        "Use of Funds:",
        "• Product (40%) — App, API integrations",
        "• Operations (30%) — Techs, training, equipment",
        "• Sales (20%) — Enterprise BD",
        "• G&A (10%) — Legal, insurance"
    ]
)

# SLIDE 11: Closing
add_slide("The future of mobility needs infrastructure.", "AUTONOMY, MAINTAINED", is_title_slide=True)

# Save
output_path = os.path.join(os.path.dirname(__file__), "MOTRA-Pitch-Deck.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
